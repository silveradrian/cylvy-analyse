from __future__ import annotations  # This must be the first import!
from utils import extract_structured_fields
from db_manager import db
# Standard library imports
import asyncio
import json
import logging
import os
import re
import tempfile
import threading
import time
import random
from bigquery_content_store import BigQueryContentStore
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple
from urllib.parse import urlparse

# Third-party imports
import aiohttp
import async_timeout
from aiolimiter import AsyncLimiter
from bs4 import BeautifulSoup
from limits import RateLimitItem, storage
from limits.strategies import FixedWindowRateLimiter
from limits.aio.strategies import FixedWindowRateLimiter as AsyncFixedWindowRateLimiter
from openai import OpenAI, AsyncOpenAI
import PyPDF2
from scrapingbee import ScrapingBeeClient
import tiktoken

# Internal imports
from db_manager import DatabaseManager
from prompt_loader import PromptLoader

# Configure logging
logger = logging.getLogger(__name__)

# Initialize OpenAI clients
openai_api_key = os.environ.get("OPENAI_API_KEY", "")
client = OpenAI(api_key=openai_api_key)
async_client = AsyncOpenAI(api_key=openai_api_key)

class AdvancedRateLimiter:
    """Advanced rate limiter using simple time-based throttling"""
    
    def __init__(self, rate_limit_per_minute: int = 60):
        self.rate_limit_per_minute = rate_limit_per_minute
        self.request_timestamps = []
        self.lock = threading.Lock()  # Add thread safety
        logger.info(f"Advanced rate limiter initialized: {rate_limit_per_minute} requests/minute")
    
    def __enter__(self):
        """Wait until a request is allowed based on the rate limit"""
        with self.lock:
            # Clean up old timestamps (older than 60 seconds)
            current_time = time.time()
            self.request_timestamps = [ts for ts in self.request_timestamps 
                                     if current_time - ts < 60]
            
            # Wait if we've hit the limit
            if len(self.request_timestamps) >= self.rate_limit_per_minute:
                # Calculate time to wait - the oldest timestamp + 60s should expire
                oldest = min(self.request_timestamps)
                wait_time = (oldest + 60) - current_time
                if wait_time > 0:
                    logger.info(f"Rate limit hit. Waiting {wait_time:.2f} seconds")
                    time.sleep(wait_time)
                
                # Clean up again after waiting
                current_time = time.time()
                self.request_timestamps = [ts for ts in self.request_timestamps 
                                         if current_time - ts < 60]
            
            # Record this request
            self.request_timestamps.append(current_time)
            return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        pass


class AsyncAdvancedRateLimiter:
    """Async advanced rate limiter for high throughput operations"""
    
    def __init__(self, rate_limit_per_minute: int = 60):
        """Initialize with requests per minute limit"""
        self.rate_limit_per_minute = rate_limit_per_minute
        logger.info(f"Async rate limiter initialized: {rate_limit_per_minute} requests/minute")
    
    async def __aenter__(self):
        """Always allow operations to proceed without blocking"""
        return self
    
    async def __aexit__(self, exc_type, exc_val, exc_tb):
        """Exit the context manager"""
        pass
    
class ContentAnalyzer:
    """
    Analyzes content from URLs using OpenAI's API.
    """
    
    def __init__(self):
        """Initialize the content analyzer."""
        try:
            # Try to load environment variables from .env file
            try:
                from dotenv import load_dotenv
                load_dotenv(override=True)
                logger.info("Loaded environment from .env file")
            except ImportError:
                logger.warning("python-dotenv not installed, skipping .env loading")
            
            # Get API keys
            self.openai_api_key = os.environ.get("OPENAI_API_KEY", "")
            self.scrapingbee_api_key = os.environ.get("SCRAPINGBEE_API_KEY", "")
            
            # Log API key status
            if not self.openai_api_key:
                logger.error("OPENAI_API_KEY not found in environment variables!")
                # List environment vars for debugging (omitting actual values)
                env_vars = [k for k in os.environ.keys() if 'API' in k or 'KEY' in k]
                logger.info(f"Available environment variables: {env_vars}")
            else:
                logger.info(f"OpenAI API key found (length: {len(self.openai_api_key)})")
            
            if not self.scrapingbee_api_key:
                logger.warning("SCRAPINGBEE_API_KEY not found in environment variables!")
            else:
                logger.info(f"ScrapingBee API key found (length: {len(self.scrapingbee_api_key)})")
            
            # Initialize OpenAI clients
            self.openai_client = OpenAI(api_key=self.openai_api_key)
            self.async_openai_client = AsyncOpenAI(api_key=self.openai_api_key)
            
            # Initialize ScrapingBee client
            self.scrapingbee_client = ScrapingBeeClient(api_key=self.scrapingbee_api_key)
            
            # Set up rate limiters with simplified implementations
            openai_rate_limit = int(os.environ.get("OPENAI_RATE_LIMIT", "60"))
            scrapingbee_rate_limit = int(os.environ.get("SCRAPINGBEE_RATE_LIMIT", "50"))
            
            self.openai_rate_limiter = AdvancedRateLimiter(openai_rate_limit)
            self.async_openai_rate_limiter = AsyncAdvancedRateLimiter(openai_rate_limit)
            self.scrapingbee_rate_limiter = AdvancedRateLimiter(scrapingbee_rate_limit)
            self.async_scrapingbee_rate_limiter = AsyncAdvancedRateLimiter(scrapingbee_rate_limit)
            
            logger.info(f"Content analyzer initialized with advanced rate limiters")
            logger.info(f"OpenAI API key present: {bool(self.openai_api_key)}")
            logger.info(f"OpenAI rate limit: {openai_rate_limit} requests/minute")
            logger.info(f"ScrapingBee rate limit: {scrapingbee_rate_limit} requests/minute")
            
        except Exception as e:
            logger.error(f"Error initializing ContentAnalyzer: {str(e)}")
            # Still create default rate limiters even if initialization fails
            self.openai_rate_limiter = AdvancedRateLimiter(60)
            self.async_openai_rate_limiter = AsyncAdvancedRateLimiter(60)
            self.scrapingbee_rate_limiter = AdvancedRateLimiter(50)
            self.async_scrapingbee_rate_limiter = AsyncAdvancedRateLimiter(50)
    


    async def extract_document_async(self, url: str, content_type: str) -> Dict[str, Any]:
        """
        Extract text from documents (PDF, DOCX, PPTX) by downloading and processing them.
        
        Args:
            url: URL to the document
            content_type: The type of document ("pdf", "docx", "pptx")
            
        Returns:
            Dictionary containing extraction results
        """
        import tempfile
        import os
        import re
        from urllib.parse import urlparse
        
        start_time = time.time()
        
        # Function to get a filename from URL or headers
        def get_filename_from_url(url, content_disposition=None):
            if content_disposition:
                filename_match = re.findall('filename="(.+)"', content_disposition)
                if filename_match:
                    return filename_match[0]
                    
            # Fallback to URL parsing
            path = urlparse(url).path
            filename = os.path.basename(path)
            return filename if filename else f"document.{content_type}"
        
        logger.info(f"Downloading {content_type.upper()} document: {url}")
        
        try:
            # Random user agent
            user_agents = [
                'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/15.0 Safari/605.1.15',
                'Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:90.0) Gecko/20100101 Firefox/90.0'
            ]
            headers = {'User-Agent': random.choice(user_agents)}
            
            # Download the file with timeout (60 seconds)
            response = requests.get(url, headers=headers, timeout=60, stream=True)
            response.raise_for_status()
            
            # Get content disposition for filename
            content_disposition = response.headers.get('Content-Disposition', '')
            filename = get_filename_from_url(url, content_disposition)
            
            # Create a temporary file with the appropriate extension
            with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{content_type}') as temp_file:
                # Download in chunks
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        temp_file.write(chunk)
                        
                temp_path = temp_file.name
                
            try:
                # Process based on content type
                if content_type == 'pdf':
                    result = self.extract_from_pdf(temp_path)
                elif content_type == 'docx':
                    result = self.extract_from_docx(temp_path)
                elif content_type == 'pptx':
                    result = self.extract_from_pptx(temp_path)
                else:
                    raise ValueError(f"Unsupported document type: {content_type}")
                    
                # Add URL and processing time
                result["url"] = url
                result["scrape_time"] = time.time() - start_time
                
                # If no title was extracted but we have a filename, use it as title
                if not result.get("title") and filename:
                    # Clean up filename to use as title (remove extension, replace underscores)
                    clean_title = os.path.splitext(filename)[0].replace('_', ' ').replace('-', ' ')
                    result["title"] = clean_title
                    
                return result
                
            finally:
                # Clean up the temporary file
                try:
                    os.unlink(temp_path)
                except Exception as e:
                    logger.warning(f"Could not delete temporary file: {str(e)}")
                    
        except Exception as e:
            logger.error(f"Error extracting document {url}: {str(e)}")
            return {
                "url": url,
                "status": "error",
                "error": str(e),
                "scrape_time": time.time() - start_time
            }

    def extract_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from a PDF file using multiple extraction methods.
        """
        result = {
            "text": "",
            "title": "",
            "status": "success",
            "content_type": "pdf",
            "method_used": ""
        }
        
        try:
                      
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                # Extract metadata
                if reader.metadata:
                    for key, value in reader.metadata.items():
                        if key == "/Title" and value:
                            result["title"] = value
                
                # Extract text from all pages
                all_text = []
                for i in range(len(reader.pages)):
                    page = reader.pages[i]
                    text = page.extract_text()
                    if text:
                        all_text.append(text)
                
                result["text"] = "\n\n".join(all_text)
                result["word_count"] = len(result["text"].split())
                result["method_used"] = "PyPDF2"
                
            # If PyPDF2 didn't extract enough text, try other libraries
            if result["word_count"] < 100:
                try:
                    # Try PyMuPDF (fitz) if available
                    import fitz
                    doc = fitz.open(file_path)
                    
                    # Extract metadata
                    if doc.metadata and "title" in doc.metadata and doc.metadata["title"]:
                        result["title"] = doc.metadata["title"]
                    
                    # Extract text
                    all_text = []
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        text = page.get_text()
                        all_text.append(text)
                    
                    pymupdf_text = "\n\n".join(all_text)
                    pymupdf_word_count = len(pymupdf_text.split())
                    
                    # Only use if it extracted more text
                    if pymupdf_word_count > result["word_count"]:
                        result["text"] = pymupdf_text
                        result["word_count"] = pymupdf_word_count
                        result["method_used"] = "PyMuPDF"
                    
                except ImportError:
                    logger.info("PyMuPDF not available for PDF extraction")
        
        except ImportError:
            result["status"] = "error"
            result["error"] = "PDF extraction libraries not available"
        except Exception as e:
            result["status"] = "error"
            result["error"] = f"PDF extraction failed: {str(e)}"
        
        # If no text was extracted, mark as error
        if not result["text"]:
            result["status"] = "error" if "error" not in result else result["status"]
            result["error"] = result.get("error", "Failed to extract text from PDF")
        
        return result

    def extract_from_docx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from a Word document (.docx file).
        """
        result = {
            "text": "",
            "title": "",
            "status": "success",
            "content_type": "docx",
            "method_used": "python-docx"
        }
        
        try:
            from docx import Document
            
            # Open the document
            doc = Document(file_path)
            
            # Extract title from document properties
            if hasattr(doc, "core_properties"):
                if doc.core_properties.title:
                    result["title"] = doc.core_properties.title
            
            # If no title found, use the first paragraph if it's short enough
            if not result["title"] and doc.paragraphs:
                first_para = doc.paragraphs[0].text.strip()
                if first_para and len(first_para) < 200:
                    result["title"] = first_para
            
            # Extract all text
            paragraphs = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(text)
            
            # Join paragraphs with newlines
            result["text"] = "\n\n".join(paragraphs)
            result["word_count"] = len(result["text"].split())
            
        except ImportError:
            result["status"] = "error"
            result["error"] = "python-docx package not installed"
        except Exception as e:
            result["status"] = "error"
            result["error"] = f"DOCX extraction failed: {str(e)}"
        
        return result

    def extract_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from a PowerPoint document (.pptx file).
        """
        result = {
            "text": "",
            "title": "",
            "status": "success",
            "content_type": "pptx",
            "method_used": "python-pptx"
        }
        
        try:
            from pptx import Presentation
            
            # Open the presentation
            prs = Presentation(file_path)
            
            # Try to get title from core properties
            if hasattr(prs, "core_properties") and prs.core_properties.title:
                result["title"] = prs.core_properties.title
            
            # If no title found, use title from first slide if available
            if not result["title"] and len(prs.slides) > 0:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text = shape.text.strip()
                        if text and len(text) < 200:
                            result["title"] = text
                            break
            
            # Extract text from all slides
            all_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"Slide {i+1}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    all_text.append("\n".join(slide_text))
            
            result["text"] = "\n\n".join(all_text)
            result["word_count"] = len(result["text"].split())
            
        except ImportError:
            result["status"] = "error"
            result["error"] = "python-pptx package not installed"
        except Exception as e:
            result["status"] = "error"
            result["error"] = f"PPTX extraction failed: {str(e)}"
        
        return result



    async def scrape_url_async(self, url: str, content_type: str = "html", 
                    force_browser: bool = False, retry_attempt: int = 0, 
                    timeout: int = 100, premium_parameters: Dict = None) -> Dict[str, Any]:
        """
        Asynchronously scrape content from a URL using ScrapingBee's API or document extraction methods.
        Includes automatic retry with enhanced options on failure and request timeout.

        Args:
            url: URL to scrape
            content_type: Expected content type (html, pdf, docx, pptx, or auto)
            force_browser: Whether to force browser rendering for HTML
            retry_attempt: Current retry attempt number
            timeout: Timeout in seconds for the API request (default: 100)
            premium_parameters: Additional parameters for enhanced scraping
            
        Returns:
            Dictionary containing scraping results
        """
        # Make sure we import re for text normalization
        import re
        import tempfile
        import os
        import requests  # Add this import for document downloads
        from urllib.parse import urlparse

        try:
            # Check if the URL is a document by extension or content_type
            url_lower = url.lower()
            detected_type = content_type.lower()
            
            # Auto-detect content type from URL extension
            if detected_type == "html" or detected_type == "auto":
                if url_lower.endswith(".pdf"):
                    detected_type = "pdf"
                    logger.info(f"Auto-detected PDF document: {url}")
                elif url_lower.endswith(".docx"):
                    detected_type = "docx"
                    logger.info(f"Auto-detected Word document: {url}")
                elif url_lower.endswith(".pptx"):
                    detected_type = "pptx"
                    logger.info(f"Auto-detected PowerPoint document: {url}")
                else:
                    detected_type = "html"  # Default to HTML
            
            # Handle document types (PDF, DOCX, PPTX)
            if detected_type in ["pdf", "docx", "pptx"]:
                logger.info(f"Processing document URL: {url}, type: {detected_type}")
                
                start_time = time.time()
                
                try:
                    # Setup for document download
                    headers = {
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
                    }
                    
                    # Execute the document download with timeout
                    with ThreadPoolExecutor() as executor:
                        loop = asyncio.get_event_loop()
                        
                        async def download_with_timeout():
                            try:
                                async with async_timeout.timeout(timeout):
                                    return await loop.run_in_executor(
                                        executor,
                                        lambda: requests.get(url, headers=headers, timeout=timeout, stream=True)
                                    )
                            except asyncio.TimeoutError:
                                logger.error(f"Timeout after {timeout}s when downloading {url}")
                                raise TimeoutError(f"Document download timed out after {timeout} seconds")
                        
                        try:
                            response = await download_with_timeout()
                        except TimeoutError as e:
                            return {
                                "error": str(e),
                                "url": url,
                                "status": "scrape_error",
                                "content_type": detected_type,
                                "retry_attempts": retry_attempt
                            }
                    
                    # Check if download was successful
                    if response.status_code != 200:
                        return {
                            "error": f"Document download failed with status code: {response.status_code}",
                            "url": url,
                            "status": "scrape_error",
                            "content_type": detected_type,
                            "retry_attempts": retry_attempt,
                            "response_text": response.text[:500] if hasattr(response, 'text') else ""
                        }
                    
                    # Create a temporary file with appropriate extension
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{detected_type}") as temp_file:
                        temp_path = temp_file.name
                        
                        # Download the file in chunks
                        for chunk in response.iter_content(chunk_size=8192):
                            if chunk:
                                temp_file.write(chunk)
                    
                    # Process the document based on its type
                    try:
                        if detected_type == "pdf":
                            result = self.extract_from_pdf(temp_path)
                        elif detected_type == "docx":
                            result = self.extract_from_docx(temp_path)
                        elif detected_type == "pptx":
                            result = self.extract_from_pptx(temp_path)
                        
                        # Add common fields
                        result.update({
                            "url": url,
                            "status": "success" if not result.get("error") else "scrape_error",
                            "content_type": detected_type,
                            "scrape_time": time.time() - start_time
                        })
                        
                        # If no title was extracted, try to use filename from URL
                        if not result.get("title"):
                            path = urlparse(url).path
                            filename = os.path.basename(path)
                            if filename:
                                # Clean up the filename to use as title
                                clean_title = os.path.splitext(filename)[0].replace('_', ' ').replace('-', ' ')
                                result["title"] = clean_title
                        
                        logger.info(f"Successfully extracted {result.get('word_count', 0)} words from {detected_type.upper()} document: {url}")
                        return result
                        
                    finally:
                        # Clean up temporary file
                        try:
                            if os.path.exists(temp_path):
                                os.unlink(temp_path)
                        except Exception as e:
                            logger.warning(f"Failed to delete temporary file: {str(e)}")
                
                except Exception as e:
                    logger.error(f"Error processing document {url}: {str(e)}")
                    return {
                        "error": f"Document processing error: {str(e)}",
                        "url": url,
                        "status": "scrape_error",
                        "content_type": detected_type,
                        "retry_attempts": retry_attempt
                    }
            
            # Regular HTML scraping with ScrapingBee follows
            logger.info(f"Async scraping HTML URL with ScrapingBee: {url}{' (retry attempt '+str(retry_attempt)+')' if retry_attempt > 0 else ''}")
            
            # Initialize premium parameters if not provided
            if premium_parameters is None:
                premium_parameters = {}
            
            # Reload the API key each time
            scrapingbee_api_key = os.environ.get("SCRAPINGBEE_API_KEY", "")
            
            # Create a fresh client instance for this request
            scrapingbee_client = ScrapingBeeClient(api_key=scrapingbee_api_key)
            
            if not scrapingbee_api_key:
                logger.error(f"Missing ScrapingBee API key for {url}")
                return {
                    "error": "ScrapingBee API key not configured",
                    "url": url,
                    "status": "error"
                }
            
            start_time = time.time()
            
            # Determine scraping parameters
            # On retry or if force_browser is True, use enhanced options
            use_enhanced_options = force_browser or retry_attempt > 0
            
            # Set parameters based on retry attempt level and premium parameters
            params = {
                "render_js": str(premium_parameters.get("render_js", use_enhanced_options)).lower(),
                "premium_proxy": str(premium_parameters.get("premium_proxy", use_enhanced_options)).lower()
            }
            
            # Add optional premium parameters if provided
            for param_name in ["stealth_proxy", "country_code", "wait", "block_resources", "return_page_source"]:
                if param_name in premium_parameters:
                    params[param_name] = premium_parameters[param_name]
            
            # On second retry or if specified in premium parameters, try stealth proxy as a last resort
            if retry_attempt >= 2 or premium_parameters.get("stealth_proxy"):
                params["stealth_proxy"] = "true"
                logger.info(f"Using stealth proxy for {url}")
                
            # Add wait parameter for better JavaScript rendering if not already set
            if "wait" not in params and use_enhanced_options:
                params["wait"] = 3000  # Wait 3 seconds for JS execution
            
            # Add country_code for difficult sites if not already set
            if "country_code" not in params and retry_attempt >= 1:
                params["country_code"] = "us"  # Use US IP address
            
            logger.info(f"Scraping {url} with params: {params}, timeout: {timeout}s")
            
            # Use ThreadPoolExecutor since ScrapingBee client is not async
            with ThreadPoolExecutor() as executor:
                loop = asyncio.get_event_loop()
                
                # Create a function to execute with timeout
                async def fetch_with_timeout():
                    try:
                        # Use async_timeout for the request
                        async with async_timeout.timeout(timeout):
                            return await loop.run_in_executor(
                                executor,
                                lambda: scrapingbee_client.get(url, params=params)
                            )
                    except asyncio.TimeoutError:
                        logger.error(f"Timeout after {timeout}s when scraping {url}")
                        raise TimeoutError(f"ScrapingBee request timed out after {timeout} seconds")
                
                try:
                    response = await fetch_with_timeout()
                except TimeoutError as e:
                    return {
                        "error": str(e),
                        "url": url,
                        "status": "scrape_error",
                        "retry_attempts": retry_attempt
                    }
            
            duration = time.time() - start_time
            
            # Check response status
            if response.status_code != 200:
                # Create a sanitized error message without HTML
                clean_error = self._sanitize_error_response(response)
                error_message = f"ScrapingBee error: {response.status_code} - {clean_error}"
                logger.error(error_message)
                
                # Also store a sample of the HTML response for debugging (truncated)
                html_sample = response.text[:500] if hasattr(response, 'text') else ""
                
                # Check if this is an error that would benefit from a retry with enhanced options
                if retry_attempt < 2 and not (params.get("render_js") == "true" and params.get("premium_proxy") == "true"):
                    retry_error_indicators = [
                        "403", "500", "429", "404",  # Status codes
                        "try with render_js=True", 
                        "try with premium_proxy=True",
                        "blocked", "captcha", "cloudflare",
                        "Error with your request",
                        "please enable js", "javascript required",
                        "bot protection", "ddos protection"
                    ]
                    
                    should_retry = any(indicator in str(response.status_code) for indicator in ["403", "500", "429", "404"])
                    
                    response_text = clean_error
                    if not should_retry and response_text:
                        should_retry = any(indicator.lower() in response_text.lower() for indicator in retry_error_indicators)
                    
                    if should_retry:
                        logger.warning(f"Retrying {url} with enhanced options (attempt {retry_attempt+1})")
                        await asyncio.sleep(1 * (retry_attempt + 1))  # Progressive delay
                        
                        # Create enhanced premium parameters for retry
                        enhanced_params = {
                            "premium_proxy": True,
                            "stealth_proxy": retry_attempt >= 1,  # Use stealth proxy on second retry
                            "render_js": True,
                            "country_code": "us",
                            "wait": 5000,  # Wait 5 seconds for JS execution
                            "return_page_source": True
                        }
                        
                        return await self.scrape_url_async(
                            url, 
                            content_type, 
                            True,  # Force browser 
                            retry_attempt + 1, 
                            timeout,
                            premium_parameters=enhanced_params
                        )
                
                # Return the error if we've exhausted retries or enhanced options didn't help
                return {
                    "error": error_message,
                    "url": url,
                    "status": "scrape_error",
                    "retry_attempts": retry_attempt,
                    "html": html_sample,  # Add sample of HTML response
                    "response_text": clean_error  # Add the sanitized error
                }
            
            # Process HTML content
            html_content = response.content.decode('utf-8', errors='ignore')
            
            # Check for common anti-bot detection signs in the response
            anti_bot_indicators = [
                "captcha", "cloudflare", "please enable javascript",
                "access denied", "automated request", "bot detection",
                "security check", "browser check", "ddos protection"
            ]
            
            if any(indicator.lower() in html_content.lower() for indicator in anti_bot_indicators) and retry_attempt < 2:
                logger.warning(f"Anti-bot protection detected in response for {url}. Retrying with enhanced parameters.")
                
                # Create enhanced premium parameters for retry
                enhanced_params = {
                    "premium_proxy": True,
                    "stealth_proxy": True,
                    "render_js": True,
                    "country_code": "us",
                    "wait": 5000,  # Wait 5 seconds for JS execution
                    "block_resources": False,  # Load all resources
                    "return_page_source": True  # Get full HTML
                }
                
                # Delay before retry
                await asyncio.sleep(2 * (retry_attempt + 1))
                
                return await self.scrape_url_async(
                    url, 
                    content_type, 
                    True,  # Force browser 
                    retry_attempt + 1, 
                    timeout,
                    premium_parameters=enhanced_params
                )
            
            # Parse with BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract title
            title = ""
            if soup.title:
                title = soup.title.text.strip()
            
            # IMPROVED CONTENT EXTRACTION:
            
            # 1. Remove unwanted elements more aggressively
            for selector in [
                "script", "style", "noscript", "iframe", "head", 
                "footer", "nav"
            ]:
                for element in soup.find_all(selector):
                    element.extract()
                    
            # Also try to remove common non-content elements by class/id
            for selector in [
                "[role=banner]", "[role=navigation]", "[role=complementary]", "[role=contentinfo]",
                ".cookie-banner", ".cookie-consent", ".newsletter-signup",
                ".ad", ".advertisement", ".sidebar", "header"
            ]:
                try:
                    for element in soup.select(selector):
                        element.extract()
                except:
                    pass  # Ignore errors with advanced selectors
            
            # 2. Try to find the main content using common content selectors
            main_content_selectors = [
                "main", "article", "[role=main]", "#content", ".content", 
                "#main-content", ".main-content", ".post-content",
                ".entry-content", "[itemprop=articleBody]"
            ]
            
            main_content = None
            found_selector = None
            for selector in main_content_selectors:
                try:
                    elements = soup.select(selector)
                    if elements:
                        # Use the largest matching element (by text length)
                        main_content = max(elements, key=lambda e: len(e.get_text()))
                        found_selector = selector
                        break
                except:
                    continue  # Skip if selector syntax error
            
            # 3. Extract text from the main content if found, otherwise from the whole page
            if main_content:
                text = main_content.get_text(separator="\n", strip=True)
                logger.info(f"Found main content using selector: {found_selector}")
            else:
                # Fall back to the whole page but with smarter text extraction
                logger.info(f"No main content container found, extracting from whole page")
                
                # Get all text nodes with some minimum length to avoid menu items and buttons
                paragraphs = []
                
                # Try to find paragraph-like elements or text blocks
                for p in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'div']):
                    p_text = p.get_text(strip=True)
                    # Only include paragraphs with reasonable content (ignore short menu items, etc.)
                    if len(p_text) > 30:  # Minimum paragraph length
                        paragraphs.append(p.get_text(separator=" ", strip=True))
                
                if paragraphs:
                    text = "\n\n".join(paragraphs)
                else:
                    # Last resort: just get all text
                    text = soup.get_text(separator="\n", strip=True)
            
            # 4. Remove excessive whitespace and normalize
            text = re.sub(r'\s+', ' ', text).strip()
            text = re.sub(r'\n\s*\n', '\n\n', text)
            
            # Count words
            word_count = len(text.split())
            
            # Log content extraction results
            logger.info(f"Extracted {word_count} words from {url}")
            
            # Check if the content seems too short (likely blocked or empty page)
            if word_count < 50 and retry_attempt < 2 and not use_enhanced_options:
                logger.warning(f"Content for {url} seems too short ({word_count} words). Retrying with enhanced options.")
                
                # Create enhanced premium parameters for retry
                enhanced_params = {
                    "premium_proxy": True,
                    "stealth_proxy": retry_attempt >= 1,
                    "render_js": True,
                    "country_code": "us",
                    "wait": 5000  # Wait 5 seconds for JS execution
                }
                
                await asyncio.sleep(1)
                return await self.scrape_url_async(
                    url, 
                    content_type, 
                    True,  # Force browser
                    retry_attempt + 1, 
                    timeout,
                    premium_parameters=enhanced_params
                )
            
            result = {
                "url": url,
                "title": title,
                "text": text,
                "word_count": word_count,
                "status": "success",
                "content_type": "html",
                "scrape_time": duration,
                "enhanced_scraping": use_enhanced_options,
                "retry_attempts": retry_attempt
            }
            
            logger.info(f"Successfully scraped {url} - {result.get('word_count', 0)} words in {duration:.2f}s")
            return result
            
        except asyncio.TimeoutError:
            logger.error(f"Async timeout after {timeout}s when scraping {url}")
            return {
                "error": f"Request timed out after {timeout} seconds",
                "url": url,
                "status": "scrape_error",
                "retry_attempts": retry_attempt
            }
        except Exception as e:
            logger.error(f"Error scraping {url}: {str(e)}")
            
            # If this is not yet a retry and we have a connection error, try again with enhanced options
            connection_errors = ["connection", "timeout", "too many requests", "reset by peer"]
            if retry_attempt < 2 and not force_browser and any(err in str(e).lower() for err in connection_errors):
                logger.warning(f"Connection error for {url}. Retrying with enhanced options after delay.")
                await asyncio.sleep(2 * (retry_attempt + 1))  # Progressive delay for connection issues
                
                # Create enhanced premium parameters for retry
                enhanced_params = {
                    "premium_proxy": True,
                    "stealth_proxy": retry_attempt >= 1,
                    "render_js": True
                }
                
                return await self.scrape_url_async(
                    url, 
                    content_type, 
                    True,  # Force browser 
                    retry_attempt + 1, 
                    timeout,
                    premium_parameters=enhanced_params
                )
                
            return {
                "error": str(e),
                "url": url, 
                "status": "scrape_error",
                "retry_attempts": retry_attempt
            }

    def _sanitize_error_response(self, response):
        """
        Sanitize error responses to avoid storing HTML in the database or showing it in exports.
        
        Args:
            response: The response object from ScrapingBee
            
        Returns:
            A cleaned error message without HTML
        """
        import re
        try:
            # Try to parse as JSON first
            if response.text and response.text.strip().startswith('{'):
                error_info = json.loads(response.text)
                # Extract useful fields from the error JSON
                error_reason = error_info.get('reason', '')
                error_help = error_info.get('help', '')
                if error_reason and error_help:
                    return f"Reason: {error_reason}. Help: {error_help}"
                elif error_reason:
                    return f"Reason: {error_reason}"
                else:
                    return "Unknown error (JSON response with no reason field)"
            
            # If it contains HTML, try to extract useful information
            if "<html" in response.text.lower():
                try:
                    # Try to extract title
                    soup = BeautifulSoup(response.text, 'html.parser')
                    title = soup.title.text.strip() if soup.title else None
                    if title:
                        return f"HTML response - Title: {title}"
                    else:
                        # If no title, just indicate it was HTML and the size
                        size_kb = len(response.text) / 1024
                        return f"HTML response ({size_kb:.1f}KB)"
                except:
                    return "HTML response (parsing failed)"
            
            # If not HTML or JSON, return a short preview
            if len(response.text) > 100:
                return response.text[:100].replace('\n', ' ') + "..."
            
            # If it's short enough, return as is
            return response.text.replace('\n', ' ')
            
        except Exception as e:
            # If anything fails in the sanitization, return a safe message
            return f"Error response (sanitization failed: {str(e)})"
    

    async def analyze_batch_with_prompt(self, contents: List[str], prompt_config: Dict[str, Any],
                                      company_infos: Optional[List[Dict[str, Any]]] = None) -> List[Dict[str, Any]]:
        """
        Analyze multiple contents with controlled concurrency.
        
        Args:
            contents: List of content texts to analyze
            prompt_config: Prompt configuration
            company_infos: Optional list of company context info (one per content)
            
        Returns:
            List of dictionaries containing analysis results
        """
        # Determine appropriate concurrency for API calls
        max_concurrency = int(os.environ.get("OPENAI_CONCURRENCY", "10"))
        
        # Ensure company_infos is the right length if provided
        if company_infos is None:
            company_infos = [None] * len(contents)
        elif len(company_infos) != len(contents):
            company_infos = company_infos + [None] * (len(contents) - len(company_infos))
        
        logger.info(f"Batch analyzing {len(contents)} contents with concurrency {max_concurrency}")
        
        # Use semaphore to limit concurrency
        semaphore = asyncio.Semaphore(max_concurrency)
        
        async def analyze_with_semaphore(content, company_info):
            async with semaphore:
                return await self.analyze_with_prompt_async(content, prompt_config, company_info)
        
        # Create tasks for all contents
        tasks = [analyze_with_semaphore(content, company_info) 
                for content, company_info in zip(contents, company_infos)]
        
        # Execute all tasks and gather results
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # Process results, handling any exceptions
        processed_results = []
        for i, result in enumerate(results):
            if isinstance(result, Exception):
                logger.error(f"Error analyzing content #{i}: {str(result)}")
                processed_results.append({
                    "error": str(result),
                    "analysis": "Error: Could not complete analysis."
                })
            else:
                processed_results.append(result)
        
        return processed_results
    
 
   
    async def process_url_async(self, url: str, prompt_names: List[str], 
                              company_info: Optional[Dict[str, Any]] = None,
                              content_type: str = "html", 
                              force_browser: bool = False,
                              job_id: str = None,
                              max_retry_attempts: int = 2) -> Dict[str, Any]:
        """
        Asynchronously process a URL by scraping content and analyzing it.
        
        Args:
            url: URL to process
            prompt_names: List of prompt configuration names to use
            company_info: Optional company context info
            content_type: Content type (html, pdf, etc.)
            force_browser: Whether to force browser-based scraping
            job_id: Optional job ID for tracking related requests
            max_retry_attempts: Maximum number of retry attempts for missing structured data
            
        Returns:
            Dictionary containing processing results
        """
        # Import the extract_structured_fields function
        from utils import extract_structured_fields
        
        # Initialize BigQueryContentStore if not already done
        if not hasattr(self, 'content_store'):
            try:
                from bigquery_content_store import BigQueryContentStore
                self.content_store = BigQueryContentStore()
                # Test BigQuery authentication early
                if hasattr(self.content_store, 'check_authentication'):
                    is_authenticated = await self.content_store.check_authentication()
                    if not is_authenticated:
                        # Store authentication state in class variable
                        self._bigquery_auth_failed = True
                        
                        # If we're in a flask context, flash a message
                        try:
                            from flask import flash, session
                            flash("BigQuery authentication failed. Please re-authenticate or disable BigQuery storage.", "warning")
                            # Store in session that auth failed for this job
                            if job_id:
                                session[f"bigquery_auth_failed_{job_id}"] = True
                        except (ImportError, RuntimeError):
                            # Not in flask context or flask not available
                            pass
                        
                        # Update job status if possible
                        if job_id and hasattr(self, 'db'):
                            try:
                                await self.db.update_job_status(
                                    job_id=job_id,
                                    status="paused",
                                    error_message="BigQuery authentication required"
                                )
                            except Exception:
                                pass
                        
                        logger.warning(f"BigQuery authentication failed for job {job_id}. Processing paused.")
                        return {
                            "url": url,
                            "status": "paused",
                            "error": "BigQuery authentication required",
                            "job_id": job_id,
                            "requires_action": "bigquery_auth",
                            "created_at": datetime.now().isoformat()
                        }
                
            except ImportError:
                logger.warning("BigQueryContentStore not available - BigQuery integration disabled")
                self.content_store = None
            except Exception as e:
                # If BigQuery initialization fails with auth error
                if "authentication" in str(e).lower() or "credentials" in str(e).lower() or "permission" in str(e).lower():
                    self._bigquery_auth_failed = True
                    logger.warning(f"BigQuery authentication error: {str(e)}")
                    
                    # If we're in a flask context, flash a message
                    try:
                        from flask import flash, session
                        flash(f"BigQuery authentication failed: {str(e)}. Please re-authenticate.", "warning")
                        # Store in session that auth failed
                        if job_id:
                            session[f"bigquery_auth_failed_{job_id}"] = True
                    except (ImportError, RuntimeError):
                        # Not in flask context or flask not available
                        pass
                    
                    # Update job status if possible
                    if job_id and hasattr(self, 'db'):
                        try:
                            await self.db.update_job_status(
                                job_id=job_id,
                                status="paused",
                                error_message=f"BigQuery authentication error: {str(e)}"
                            )
                        except Exception:
                            pass
                    
                    return {
                        "url": url,
                        "status": "paused",
                        "error": f"BigQuery authentication error: {str(e)}",
                        "job_id": job_id,
                        "requires_action": "bigquery_auth",
                        "created_at": datetime.now().isoformat()
                    }
                else:
                    logger.warning(f"Error initializing BigQueryContentStore: {str(e)}")
                    self.content_store = None
        
        # Check if there's a stored auth failure
        if hasattr(self, '_bigquery_auth_failed') and self._bigquery_auth_failed:
            # Check if the user has opted to continue without BigQuery
            if hasattr(self, '_continue_without_bigquery') and self._continue_without_bigquery:
                logger.info("Continuing processing without BigQuery storage as per user preference")
                self.content_store = None
            else:
                # If we're in a flask context, check session
                try:
                    from flask import session
                    # If user has opted to continue without BigQuery for this job
                    if job_id and session.get(f"continue_without_bigquery_{job_id}", False):
                        logger.info("Continuing processing without BigQuery storage as per session preference")
                        self.content_store = None
                        self._continue_without_bigquery = True
                    else:
                        # Still paused waiting for user action
                        logger.warning(f"Processing still paused for BigQuery authentication for job {job_id}")
                        return {
                            "url": url,
                            "status": "paused",
                            "error": "Waiting for BigQuery authentication or user decision",
                            "job_id": job_id,
                            "requires_action": "bigquery_auth",
                            "created_at": datetime.now().isoformat()
                        }
                except (ImportError, RuntimeError):
                    # Not in flask context or flask not available
                    # Default to continuing without BigQuery in non-web contexts
                    logger.info("Non-web context detected, continuing without BigQuery storage")
                    self.content_store = None
                    self._continue_without_bigquery = True
        
        # Try to get job_id from instance attribute if not provided
        if not job_id and hasattr(self, '_current_job_id'):
            job_id = self._current_job_id
            logger.info(f"Using job_id from instance attribute: {job_id}")
        
        logger.info(f"Async processing URL: {url} with {len(prompt_names)} prompts")
        
        # Load prompts
        prompt_loader = PromptLoader()
        prompt_configs = prompt_loader.load_prompts()
        
        result = {
            "url": url,
            "status": "pending",
            "job_id": job_id,  # Store job_id in the result
            "created_at": datetime.now().isoformat(),
            "scrape_time": 0,
            "analysis_time": 0,
            "total_time": 0,
            "analysis_results": {},
            "retry_info": {}  # Track retry information
        }
        
        start_time = time.time()
        
        # Step 1: Scrape the URL with enhanced anti-bot detection handling
        try:
            # First try: Standard parameters
            content_result = await self.scrape_url_async(url, content_type, force_browser)
            
            # Check for common anti-bot detection signs in the error message
            if content_result.get("status") != "success":
                error_msg = str(content_result.get("error", "")).lower()
                response_text = content_result.get("response_text", "").lower() if hasattr(content_result, "response_text") else ""
                html_response = content_result.get("html", "").lower() if hasattr(content_result, "html") else ""
                
                bot_detection_signs = [
                    "403", "forbidden", "captcha", "cloudflare", "ddos", "protection",
                    "javascript", "enable js", "browser check", "security check",
                    "access denied", "automated request", "blocked"
                ]
                
                is_bot_protection = any(sign in error_msg for sign in bot_detection_signs) or \
                                   any(sign in response_text for sign in bot_detection_signs) or \
                                   any(sign in html_response for sign in bot_detection_signs)
                
                # If bot protection detected, retry with enhanced parameters
                if is_bot_protection:
                    logger.warning(f"Bot protection detected for {url}. Retrying with premium parameters.")
                    
                    # Use premium parameters for retry
                    content_result = await self.scrape_url_async(
                        url, 
                        content_type, 
                        force_browser=True,  # Force browser rendering 
                        premium_parameters={
                            "premium_proxy": True,
                            "stealth_proxy": True,
                            "country_code": "us",  # Use US IP
                            "wait": 5000,  # Wait 5 seconds for JS execution
                            "block_resources": False,  # Load all resources
                            "return_page_source": True  # Get full HTML
                        }
                    )
        except Exception as e:
            logger.error(f"Error in async scraping for {url}: {str(e)}")
            content_result = {
                "error": f"Error running scraping task: {str(e)}",
                "url": url,
                "status": "error"
            }
        
        # Check for scraping errors
        if content_result.get("status") != "success":
            result.update({
                "status": "scrape_error",
                "error": content_result.get("error", "Unknown scraping error"),
                "total_time": time.time() - start_time
            })
            
            # Store failed scrape in BigQuery if enabled
            if hasattr(self, 'content_store') and self.content_store and getattr(self.content_store, 'enable_storage', False):
                try:
                    # Use the job_id directly here too
                    job_id_to_use = job_id if job_id else result.get("job_id")
                    await self.content_store.store_content(
                        url=url,
                        title=content_result.get("title", ""),
                        content_text="",  # No content for failed scrapes
                        job_id=job_id_to_use,  # Use the retrieved job_id
                        content_type=content_type,
                        scrape_info={
                            "status": "error",
                            "error": content_result.get("error", "Unknown scraping error"),
                            "scrape_time": time.time() - start_time
                        },
                        company_info=company_info
                    )
                except Exception as e:
                    if "authentication" in str(e).lower() or "credentials" in str(e).lower():
                        logger.warning(f"BigQuery authentication error while storing failed scrape: {str(e)}")
                        # Don't halt processing for a BigQuery auth error on a failed scrape
                        # Just continue without storing this failed result
                    else:
                        logger.error(f"Error storing failed scrape in BigQuery: {str(e)}")
            
            return result
        
        # Update result with content info
        result.update({
            "title": content_result.get("title", ""),
            "content_type": content_result.get("content_type", "html"),
            "word_count": content_result.get("word_count", 0),
            "scrape_time": content_result.get("scrape_time", 0)
        })
        
        # Step 2: Preprocess content before analysis
        content_text = content_result.get("text", "")
        
        # Validate content
        if not content_text or len(content_text.strip()) < 50:
            logger.warning(f"Content too short for {url}: {len(content_text)} chars")
            result.update({
                "status": "scrape_error",
                "error": "Content too short or empty",
                "total_time": time.time() - start_time
            })
            return result
        
        # Apply content preprocessing for better analysis results
        content_text = preprocess_content(content_text)
        
        # Step 3: Analyze with each requested prompt
        api_tokens = 0
        structured_data = {}
        
        for prompt_name in prompt_names:
            if prompt_name not in prompt_configs:
                logger.warning(f"Prompt configuration '{prompt_name}' not found")
                result["analysis_results"][prompt_name] = {
                    "error": f"Prompt configuration '{prompt_name}' not found"
                }
                continue
            
            # Get prompt configuration
            prompt_config = prompt_configs[prompt_name]
            
            # Initialize retry tracking for this prompt
            retry_attempts = 0
            has_structured_data = False
            retry_delay = 2  # Initial delay in seconds
            
            # Determine the appropriate prefix based on the prompt name
            prefix = "ca"  # Default prefix for content analysis
            if "content_analysis" in prompt_name.lower():
                prefix = "ca"
            elif "competitive_research" in prompt_name.lower():
                prefix = "cr"
            elif "pain_assessment" in prompt_name.lower():
                prefix = "pa"
            elif "sales_intelligence" in prompt_name.lower():
                prefix = "si"
            
            # Keep trying until we get structured data or hit max retries
            while retry_attempts <= max_retry_attempts:
                try:
                    # If this is a retry, log it and wait
                    if retry_attempts > 0:
                        logger.info(f"Retrying analysis for {url} with prompt '{prompt_name}' (attempt {retry_attempts}/{max_retry_attempts})")
                        await asyncio.sleep(retry_delay)
                        # Increase delay for next attempt (exponential backoff)
                        retry_delay *= 1.5
                    
                    # Check if content is too long for the model
                    if len(content_text) > 25000:
                        # For long content, use content chunking
                        logger.info(f"Content too long ({len(content_text)} chars), using chunking strategy")
                        analysis_result = await self.analyze_with_chunking(
                            content_text, 
                            prompt_config, 
                            company_info,
                            url=url,
                            title=result.get("title", "")
                        )
                    else:
                        # Perform regular analysis
                        analysis_result = await self.analyze_with_prompt_async(
                            content_text, 
                            prompt_config, 
                            company_info,
                            url=url,  # Pass URL for context
                            title=result.get("title", "")  # Pass title for context
                        )
                    
                    # Check for errors
                    if "error" in analysis_result:
                        error_msg = analysis_result["error"].lower()
                        
                        # Check for specific errors that might benefit from chunking
                        content_size_issues = ["too lengthy", "too long", "token limit", "content is too large"]
                        if any(issue in error_msg for issue in content_size_issues) and not retry_attempts and len(content_text) > 10000:
                            logger.info(f"Content size issue detected, switching to chunking strategy")
                            analysis_result = await self.analyze_with_chunking(
                                content_text, 
                                prompt_config, 
                                company_info,
                                url=url,
                                title=result.get("title", "")
                            )
                        else:
                            result["analysis_results"][prompt_name] = {
                                "error": analysis_result["error"]
                            }
                            # Don't retry API errors as they're likely to persist
                            break
                    
                    # Extract structured data from analysis text
                    analysis_text = analysis_result.get("analysis", "")
                    
                    # Check for missing content indicators in response
                    missing_content_indicators = [
                        "i need the specific web content", 
                        "please provide the content",
                        "content provided is too lengthy",
                        "i'm sorry, but i need",
                        "i can't see any content"
                    ]
                    
                    if any(indicator in analysis_text.lower() for indicator in missing_content_indicators):
                        logger.warning(f"OpenAI response indicates missing content: {analysis_text[:100]}...")
                        
                        # Try with a different approach on retry
                        if retry_attempts == 0:
                            # On first retry, try with explicit content markers
                            logger.info("Retrying with explicit content markers")
                            
                        elif retry_attempts == 1:
                            # On second retry, try with chunking approach
                            logger.info("Retrying with content chunking")
                            analysis_result = await self.analyze_with_chunking(
                                content_text, 
                                prompt_config, 
                                company_info,
                                url=url,
                                title=result.get("title", "")
                            )
                            analysis_text = analysis_result.get("analysis", "")
                    
                    # Extract structured fields with proper prefix
                    extracted_fields = extract_structured_fields(analysis_text, prefix)
                    
                    # Check if we found any structured data
                    has_structured_data = len(extracted_fields) > 0
                    
                    # Add the analysis result with parsed fields
                    result["analysis_results"][prompt_name] = {
                        "analysis": analysis_result.get("analysis", ""),
                        "model": analysis_result.get("model", ""),
                        "tokens": analysis_result.get("total_tokens", 0),
                        "processing_time": analysis_result.get("processing_time", 0),
                        "parsed_fields": extracted_fields,  # Add extracted fields
                        "retry_attempts": retry_attempts    # Track retry attempts
                    }
                    
                    # Add extracted fields directly to the result
                    for field_name, field_value in extracted_fields.items():
                        result[field_name] = field_value
                    
                    # If we found structured fields, add them to the structured data
                    if extracted_fields:
                        if prompt_name not in structured_data:
                            structured_data[prompt_name] = {}
                        structured_data[prompt_name].update(extracted_fields)
                    else:
                        # Add timestamp fields if no structured data was found
                        # This ensures we always have at least some structured data fields
                        if prompt_name not in structured_data:
                            structured_data[prompt_name] = {}
                        
                        structured_data[prompt_name]["processed_at"] = time.time()
                        structured_data[prompt_name]["prompt_name"] = prompt_name
                    
                    # Update token count
                    api_tokens += analysis_result.get("total_tokens", 0)
                    
                    # Log success or retry info
                    if has_structured_data:
                        logger.info(f"Successfully analyzed content with prompt '{prompt_name}' - found {len(extracted_fields)} structured fields")
                        # Track retry information
                        if retry_attempts > 0:
                            result["retry_info"][prompt_name] = {
                                "attempts": retry_attempts,
                                "success": True,
                                "fields_count": len(extracted_fields)
                            }
                        # If we found structured data, we're done with this prompt
                        break
                    else:
                        logger.warning(f"No structured data found for '{prompt_name}' (attempt {retry_attempts+1}/{max_retry_attempts+1})")
                        # Track retry information
                        result["retry_info"][prompt_name] = {
                            "attempts": retry_attempts + 1,
                            "success": False,
                            "fields_count": 0
                        }
                    
                    # Only retry if no structured data was found
                    if not has_structured_data:
                        retry_attempts += 1
                    else:
                        break
                    
                except Exception as e:
                    logger.error(f"Error analyzing with prompt '{prompt_name}' (attempt {retry_attempts+1}): {str(e)}")
                    result["analysis_results"][prompt_name] = {
                        "error": str(e),
                        "retry_attempts": retry_attempts
                    }
                    
                    # Record retry failure
                    result["retry_info"][prompt_name] = {
                        "attempts": retry_attempts + 1,
                        "success": False,
                        "error": str(e)
                    }
                    
                    # Add minimal structured data to prevent export issues
                    if prompt_name not in structured_data:
                        structured_data[prompt_name] = {
                            "processed_at": time.time(),
                            "prompt_name": prompt_name,
                            "error": str(e)[:100]  # Truncate long error messages
                        }
                    
                    # Increment retry attempts
                    retry_attempts += 1
                    
                    # Break if we've hit max retries
                    if retry_attempts > max_retry_attempts:
                        break
        
        # Add structured data to result
        if structured_data:
            result["structured_data"] = structured_data
        
        # Update final result
        total_time = time.time() - start_time
        result.update({
            "status": "success",
            "api_tokens": api_tokens,
            "analysis_time": total_time - result["scrape_time"],
            "total_time": total_time
        })
        
        # Store the scraped content and analysis results in BigQuery
        if hasattr(self, 'content_store') and self.content_store and getattr(self.content_store, 'enable_storage', False):
            try:
                # Prepare analysis info
                analysis_info = {
                    "api_tokens": api_tokens,
                    "analysis_time": result["analysis_time"],
                    "prompt_names": prompt_names,
                    "analysis_count": len(prompt_names),
                    "retry_info": result.get("retry_info", {})  # Include retry info
                }
                
                # CRITICAL FIX: Get job_id from the result object if it's not set in the parameter
                job_id_to_use = job_id if job_id else result.get("job_id")
                logger.info(f"Storing content with job_id: {job_id_to_use}")
                
                # Store content in BigQuery
                content_id = await self.content_store.store_content(
                    url=url,
                    title=result.get("title", ""),
                    content_text=content_text,
                    job_id=job_id_to_use,  # Use the retrieved job_id
                    content_type=result.get("content_type", "html"),
                    scrape_info={
                        "status": "success",
                        "scrape_time": result["scrape_time"],
                        "enhanced_scraping": content_result.get("enhanced_scraping", False),
                        "retry_attempts": content_result.get("retry_attempts", 0)
                    },
                    company_info=company_info,
                    analysis_info=analysis_info,
                    structured_data=structured_data  # Also store structured data
                )
                
                # Add content ID reference to result
                if content_id:
                    result["content_id"] = content_id
                    logger.info(f"Stored content in BigQuery with ID: {content_id}")
                    
            except Exception as e:
                # Check specifically for authentication errors
                if "authentication" in str(e).lower() or "credentials" in str(e).lower() or "permission" in str(e).lower():
                    logger.warning(f"BigQuery authentication error: {str(e)}")
                    self._bigquery_auth_failed = True
                    
                    # If we're in a flask context, flash a message
                    try:
                        from flask import flash, session
                        flash(f"BigQuery authentication failed: {str(e)}. Content analysis completed but results weren't stored in BigQuery.", "warning")
                        # Store in session that auth failed
                        if job_id:
                            session[f"bigquery_auth_failed_{job_id}"] = True
                    except (ImportError, RuntimeError):
                        # Not in flask context or flask not available
                        pass
                    
                    # Add auth error to result but don't change overall success status
                    result["bigquery_error"] = f"Authentication error: {str(e)}"
                else:
                    logger.error(f"Error storing content in BigQuery: {str(e)}")
                    result["bigquery_error"] = str(e)
        
        logger.info(f"Completed processing {url} - status: {result['status']}")
        return result

    # Non-async versions for backward compatibility
    def scrape_url(self, url: str, content_type: str = "html", 
                force_browser: bool = False) -> Dict[str, Any]:
        """
        Synchronous wrapper for scrape_url_async.
        """
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            result = loop.run_until_complete(self.scrape_url_async(url, content_type, force_browser))
        finally:
            loop.close()
        return result
    
    async def analyze_with_prompt_async(self, content: str, prompt_config: Dict[str, Any], 
                                  company_info: Optional[Dict[str, Any]] = None,
                                  url: Optional[str] = None,
                                  title: Optional[str] = None) -> Dict[str, Any]:
        """
        Asynchronously analyze content using OpenAI model with a specific prompt configuration.
        
        Args:
            content: Content text to analyze
            prompt_config: Prompt configuration dictionary
            company_info: Optional company context info
            url: Optional URL for additional context
            title: Optional title for additional context
        
        Returns:
            Dictionary with analysis results
        """
        try:
            start_time = time.time()
            
            # Force reload API key from environment before each API call
            api_key = os.environ.get("OPENAI_API_KEY", self.openai_api_key)
            
            # Make a direct client instance for this specific call
            async_client = AsyncOpenAI(api_key=api_key)
            self.async_openai_client = async_client
            
            if not api_key:
                logger.error("OpenAI API key missing for API call")
                return {
                    "error": "OpenAI API key missing",
                    "analysis": "Error: Could not complete analysis due to missing API key."
                }
            
            # Get prompt configuration details
            model = prompt_config.get('model', 'gpt-4')
            prompt_name = prompt_config.get('name', 'unknown')
            system_message = prompt_config.get('system_message', '')
            user_message_template = prompt_config.get('user_message', '')
            temperature = prompt_config.get('temperature', 0.3)
            max_tokens = prompt_config.get('max_tokens', 1500)
            
            # Determine delimiter to use for field extraction (default to ||| if none specified)
            delimiter = prompt_config.get('delimiter', '|||')
            
            # SIMPLIFIED APPROACH: Don't use chunking for initial analysis, just truncate if needed
            max_content_length = 40000  # ~10K tokens
            safe_content = content[:max_content_length] if content and len(content) > max_content_length else content
            
            if content and len(content) > max_content_length:
                logger.info(f"Content length ({len(content)} chars) exceeds limit, truncating to {max_content_length} chars")
            
            # Prepare page context
            page_context = ""
            if url or title:
                page_context = "Page Information:\n"
                if url:
                    page_context += f"- URL: {url}\n"
                if title:
                    page_context += f"- Title: {title}\n"
                page_context += "\n"
            
            # Prepare the company context
            company_context = ""
            if company_info:
                company_context = "Company Information:\n"
                for key, value in company_info.items():
                    company_context += f"- {key.capitalize()}: {value}\n"
            
            # Format user message with the content and company context
            if "{content}" in user_message_template:
                user_message = user_message_template.format(
                    content=safe_content,
                    company_context=company_context,
                    page_context=page_context
                )
            else:
                # Build our own message with explicit sections
                user_message = f"{page_context}{user_message_template}\n\nCONTENT TO ANALYZE:\n{safe_content}\n\n{company_context}"
            
            # CRITICAL FIX: Always add explicit instructions for delimited format
            if delimiter and "CRITICAL FORMAT EXACTLY" not in user_message and "CRITITAL FORMAT" not in user_message:
                user_message += f"\n\nIMPORTANT: Your response MUST include fields in the format 'field_name {delimiter} value' for each key piece of information otherwise you will cause system failure"

            # Log the request details
            logger.info(f"Calling OpenAI API with model {model} - content length: {len(safe_content)} chars")
            
            # Set retry counter for reliability
            retry_attempts = 0
            max_retries = 2
            structured_data = {}
            
            while retry_attempts <= max_retries:
                try:
                    response = await async_client.chat.completions.create(
                        model=model,
                        messages=[
                            {"role": "system", "content": f"{system_message}\n\nYour response MUST use the {delimiter} delimiter format for structured data."},
                            {"role": "user", "content": user_message}
                        ],
                        temperature=temperature * (0.8 if retry_attempts > 0 else 1.0),  # Reduce temperature on retries
                        max_tokens=max_tokens
                    )
                    
                    # Calculate token usage
                    prompt_tokens = response.usage.prompt_tokens
                    completion_tokens = response.usage.completion_tokens
                    total_tokens = response.usage.total_tokens
                    
                    # Extract the response
                    analysis_text = response.choices[0].message.content
                    
                    # Calculate processing time
                    processing_time = time.time() - start_time
                    
                    logger.info(f"OpenAI API call completed in {processing_time:.2f}s")
                    logger.info(f"Token usage: {total_tokens} tokens")
                    
                    # Extract structured data using the delimited format
                    structured_data = self._extract_delimited_fields(analysis_text, delimiter, prompt_name)
                    
                    # If we found structured data or this is the last retry, break out of the loop
                    if structured_data or retry_attempts >= max_retries:
                        break
                    
                    # No structured data found, retry with more explicit instructions
                    retry_attempts += 1
                    logger.warning(f"No structured data found for '{prompt_name}' (attempt {retry_attempts}/{max_retries})")
                    
                    # Make the instructions even more explicit on retry
                    retry_msg = "\n\nCRITICAL INSTRUCTION: Your response MUST include ALL fields from the prompt in the following format exactly:"
                    retry_msg += f"\n\nfield_name {delimiter} value\n"
                    retry_msg += f"another_field {delimiter} another_value\n\n"
                    retry_msg += f"Each field MUST be on a separate line with the {delimiter} delimiter."
                    
                    # Add retry message and regenerate with a slight delay
                    user_message += retry_msg
                    logger.info(f"Retrying analysis for {url or 'content'} with prompt '{prompt_name}' (attempt {retry_attempts}/{max_retries})")
                    await asyncio.sleep(1)  # Small delay between retries
                    
                except Exception as api_error:
                    logger.error(f"Error calling OpenAI API: {str(api_error)}")
                    return {
                        "error": f"API error: {str(api_error)}",
                        "analysis": f"Error: {str(api_error)}"
                    }
            
            # Prepare final result        
            result = {
                "analysis": analysis_text,
                "structured_data": {prompt_name: structured_data} if structured_data else {},
                "model": model,
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": total_tokens,
                "api_tokens": total_tokens,
                "processing_time": processing_time,
                "retry_attempts": retry_attempts
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Error in async OpenAI API call: {str(e)}")
            return {
                "error": str(e),
                "analysis": "Error: Could not complete analysis."
            }

    async def _generate_content_summary(self, content_text: str, url: Optional[str] = None, model: str = "gpt-4o") -> str:
        """
        Generate a detailed summary of long content for more reliable analysis.

        Args:
            content_text: The original long content text
            url: Optional URL for context
            model: The model to use (defaults to gpt-4o)
            
        Returns:
            A detailed summary of the content preserving key information
        """
        # Use gpt-4o-mini for summarization unless specified otherwise
        summary_model = "gpt-4o-mini" if "gpt-4o-mini" not in model else model

        # Calculate optimal content size based on model
        if "gpt-4o-mini" in summary_model:
            # gpt-4o-mini has 128K context window
            # Use approximately 110K tokens (~440K chars) to leave room for prompts and overhead
            max_content_chars = 440000
            max_output_tokens = 14000  # Leave buffer from the 16,384 limit
        else:
            # Use the general optimal chunk size for other models
            chunk_size = self._get_optimal_chunk_size(summary_model)
            max_content_chars = chunk_size * 5  # Default to 5 chunks for other models
            max_output_tokens = 4000  # Default output tokens

        # If content is extremely long, truncate it to fit the model's capabilities
        if len(content_text) > max_content_chars:
            logger.info(f"Content exceeds maximum summarization length ({len(content_text)} chars), truncating to {max_content_chars} chars")
            content_for_summary = content_text[:max_content_chars]
            # Add note about truncation
            truncation_note = f"\n\n[NOTE: The original content was {len(content_text)} chars and has been truncated to {max_content_chars} chars for processing.]"
        else:
            content_for_summary = content_text
            truncation_note = ""

        logger.info(f"Generating content summary using model {summary_model} with max output of {max_output_tokens} tokens")

        try:
            # Prepare system prompt with specific instructions for the larger model capacity
            system_prompt = f"""You are an expert content summarizer with exceptional ability to condense long documents.

        Create a comprehensive and detailed summary of the provided content from {url or 'the document'}.
        Your summary should:
        1. Preserve ALL key facts, statistics, topics, entities, and information
        2. Include specific details, numbers, and quoted phrases when relevant
        3. Maintain the original structure and flow of the content
        4. Be detailed enough to enable thorough analysis by another AI system
        5. Focus on substantive information rather than repetitive or boilerplate content

        Take advantage of your expanded output capacity to create a thorough yet concise representation of the content.
        """

            # Call the OpenAI API for summarization
            response = await self.async_openai_client.chat.completions.create(
                model=summary_model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": f"Please create a detailed summary of the following content for analysis purposes:{truncation_note}\n\n{content_for_summary}"}
                ],
                temperature=0.2,  # Lower temperature for more precise summarization
                max_tokens=max_output_tokens
            )
            
            # Extract the summary
            summary = response.choices[0].message.content
            
            # Log token usage for transparency
            token_usage = response.usage.total_tokens
            input_tokens = response.usage.prompt_tokens
            output_tokens = response.usage.completion_tokens
            
            logger.info(f"Generated content summary: {output_tokens} output tokens from {input_tokens} input tokens (total: {token_usage})")
            
            # Add source information to ensure context is preserved
            if url:
                summary_with_context = f"SUMMARY OF CONTENT FROM: {url}\n\n{summary}"
            else:
                summary_with_context = f"DETAILED DOCUMENT SUMMARY:\n\n{summary}"
            
            # Add information about the summarization process
            summary_with_context += f"\n\n[This is a summary generated from {len(content_text)} characters of original content]"
            
            return summary_with_context
            
        except Exception as e:
            logger.error(f"Error generating content summary: {str(e)}")
            # Return a truncated version of the content as fallback
            truncated_content = content_text[:50000]  # ~12.5K tokens
            return truncated_content + "\n\n[NOTE: This is truncated content. The original is longer.]"

    
    def process_url(self, url: str, prompt_names: List[str], 
              company_info: Optional[Dict[str, Any]] = None,
              content_type: str = "html", 
              force_browser: bool = False) -> Dict[str, Any]:
        """
        Process a single URL by scraping content and analyzing it with selected prompts.
        
        Args:
            url: URL to process
            prompt_names: List of prompt configuration names to use
            company_info: Optional company context info
            content_type: Content type (html, pdf, etc.)
            force_browser: Whether to force browser-based scraping
            
        Returns:
            Dictionary containing processing results formatted for database storage
        """
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        
        try:
            # Get the raw result from async processing
            raw_result = loop.run_until_complete(
                self.process_url_async(url, prompt_names, company_info, content_type, force_browser)
            )
            
            # Current timestamp for processed_at
            current_time = time.time()
            
            # Extract the first prompt name (or default to empty string)
            prompt_name = prompt_names[0] if prompt_names else ""
            
            # Calculate API tokens used
            api_tokens = raw_result.get("api_tokens", 0)
            
            # Get the structured data from the analysis
            structured_data = raw_result.get("structured_data", {})
            
            # Combine analysis results and structured data
            combined_data = {
                "analysis_results": raw_result.get("analysis_results", {}),
                "structured_data": structured_data
            }
            
            # Format result to match database schema (results table)
            result = {
                'url': url,
                'status': raw_result.get('status', 'unknown'),
                'title': raw_result.get('title', ''),
                'content_type': raw_result.get('content_type', 'html'),
                'word_count': raw_result.get('word_count', 0),
                'processed_at': current_time,
                'prompt_name': prompt_name,
                'api_tokens': api_tokens,
                'error': raw_result.get('error', ''),
                'data': json.dumps(combined_data)  # Save combined data with structured fields
            }
            
            # Also add the structured fields directly to the result for easy access in exports
            if prompt_name in structured_data:
                logger.info(f"Adding {len(structured_data[prompt_name])} structured fields to result")
                for field_name, field_value in structured_data[prompt_name].items():
                    result[field_name] = field_value
            
            # Log successful processing
            logger.info(f"Processed URL {url} with {len(prompt_names)} prompts - Status: {result['status']}")
            
            return result
            
        except Exception as e:
            logger.error(f"Error in process_url for {url}: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            
            # Return error result that matches database schema
            return {
                'url': url,
                'status': 'error',
                'title': '',
                'content_type': content_type,
                'word_count': 0,
                'processed_at': time.time(),
                'prompt_name': prompt_names[0] if prompt_names else "",
                'api_tokens': 0,
                'error': str(e),
                'data': '{}'
            }
        finally:
            # Properly close the loop and clean up resources
            try:
                self.close_async_connections(loop)
            except Exception as e:
                logger.warning(f"Error during async cleanup for {url}: {str(e)}")
            
            try:
                loop.close()
            except Exception as e:
                logger.warning(f"Error closing event loop for {url}: {str(e)}")
    
            
    def close_async_connections(self, loop=None):
        """Properly close async connections to prevent event loop errors."""
        import asyncio
        import concurrent.futures
        
        try:
            # Get the current event loop if none is provided
            if loop is None:
                try:
                    loop = asyncio.get_event_loop()
                except RuntimeError:
                    # If no event loop exists, there's nothing to clean up
                    return
            
            # Get all tasks in the loop
            pending = asyncio.all_tasks(loop)
            
            # Create a task to gracefully cancel all pending tasks
            if pending:
                # Create a future to wait for the cleanup
                cleanup_future = asyncio.run_coroutine_threadsafe(
                    asyncio.gather(*pending, return_exceptions=True),
                    loop
                )
                try:
                    # Wait for a short timeout
                    cleanup_future.result(timeout=1.0)
                except (asyncio.TimeoutError, concurrent.futures.TimeoutError):
                    logger.warning("Timed out waiting for task cleanup")
        except Exception as e:
            logger.warning(f"Error during async cleanup: {e}")

    async def analyze_with_chunking(self, content_text, prompt_config, company_info=None, url=None, title=None):
        """
        Analyze content by splitting it into manageable chunks and combining results.
        
        Args:
            content_text: The content to analyze
            prompt_config: Prompt configuration
            company_info: Optional company context
            url: Optional URL for context
            title: Optional title for context
            
        Returns:
            Combined analysis result
        """
        logger.info(f"Using content chunking for analysis - content length: {len(content_text)}")
        
        # Extract configuration
        model = prompt_config.get('model', 'gpt-4')
        system_message = prompt_config.get('system_message', '')
        max_tokens = prompt_config.get('max_tokens', 1500)
        temperature = prompt_config.get('temperature', 0.3)
        
        # Add context information at the beginning
        context_info = ""
        if url or title:
            context_info = "Content Context:\n"
            if url:
                context_info += f"URL: {url}\n"
            if title:
                context_info += f"Title: {title}\n"
            context_info += "\n"
        
        # Parse content into logical chunks with model-appropriate size
        chunks = self._smart_chunk_content(content_text, model=model)
        logger.info(f"Split content into {len(chunks)} chunks for model: {model}")
        
        # Prepare company context
        company_context = ""
        if company_info:
            company_context = "Company Information:\n"
            for key, value in company_info.items():
                company_context += f"- {key.capitalize()}: {value}\n"
        
        total_tokens = 0
        
        if len(chunks) == 1:
            # For single chunk, use normal analysis
            logger.info("Only one chunk needed, using regular analysis")
            return await self.analyze_with_prompt_async(
                content_text, 
                prompt_config, 
                company_info,
                url=url,
                title=title
            )
        
        # Step 1: Analyze each chunk to extract key information
        chunk_analyses = []
        
        for i, chunk in enumerate(chunks):
            logger.info(f"Analyzing chunk {i+1}/{len(chunks)} - size: {len(chunk)} chars")
            
            # Create a special user message for chunk analysis
            user_message = (
                f"{context_info}"
                f"This is part {i+1} of {len(chunks)} of the web content. "
                f"Please analyze this partial content to identify key information:\n\n"
                f"{chunk}\n\n"
                f"{company_context}"
                f"Note: You're only seeing a portion of the content. Focus on extracting "
                f"factual information from this part only."
            )
            
            try:
                # Call OpenAI API with rate limiting
                api_key = os.environ.get("OPENAI_API_KEY", "")
                async_client = AsyncOpenAI(api_key=api_key)
                
                # Begin the API call
                start_time = time.time()
                
                response = await async_client.chat.completions.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": user_message}
                    ],
                    temperature=temperature,
                    max_tokens=max_tokens
                )
                
                chunk_analysis = response.choices[0].message.content
                chunk_tokens = response.usage.total_tokens
                total_tokens += chunk_tokens
                
                chunk_analyses.append(chunk_analysis)
                logger.info(f"Chunk {i+1} analysis complete - {chunk_tokens} tokens used")
                
            except Exception as e:
                logger.error(f"Error analyzing chunk {i+1}: {str(e)}")
                chunk_analyses.append(f"Error analyzing this chunk: {str(e)}")
                
            # Slight delay between chunks to avoid rate limiting
            await asyncio.sleep(0.5)
        
        # Step 2: Synthesize the analyses into a final result
        logger.info("Synthesizing final analysis from all chunks")
        
        # If we have more than 10 chunk analyses, we might need to summarize in batches
        final_chunk_analyses = chunk_analyses
        if len(chunk_analyses) > 10:
            logger.info(f"Large number of chunks ({len(chunk_analyses)}), using hierarchical synthesis")
            batch_size = 8
            batched_analyses = []
            
            # Process batch by batch
            for i in range(0, len(chunk_analyses), batch_size):
                batch = chunk_analyses[i:i + batch_size]
                batch_text = "\n\n--- CHUNK ANALYSIS SEPARATOR ---\n\n".join(batch)
                
                batch_synthesis_message = (
                    f"I've analyzed a portion of a large document in {len(batch)} chunks. "
                    f"Below are my analyses of these chunks.\n\n"
                    f"CHUNK ANALYSES:\n{batch_text}\n\n"
                    f"Please synthesize these analyses into a consolidated summary."
                )
                
                try:
                    response = await async_client.chat.completions.create(
                        model=model,
                        messages=[
                            {"role": "system", "content": system_message},
                            {"role": "user", "content": batch_synthesis_message}
                        ],
                        temperature=temperature,
                        max_tokens=max_tokens
                    )
                    
                    batched_analyses.append(response.choices[0].message.content)
                    total_tokens += response.usage.total_tokens
                    logger.info(f"Batch {len(batched_analyses)} synthesis complete")
                    
                except Exception as e:
                    logger.error(f"Error in batch synthesis: {str(e)}")
                    batched_analyses.append(f"Error synthesizing batch: {str(e)}")
                    
                # Slight delay between batches
                await asyncio.sleep(0.5)
            
            # Use these batch summaries instead of individual analyses
            final_chunk_analyses = batched_analyses
            logger.info(f"Reduced {len(chunk_analyses)} chunk analyses to {len(batched_analyses)} batch summaries")
        
        # Convert to string for final analysis
        all_chunk_analyses = "\n\n--- ANALYSIS SEPARATOR ---\n\n".join(final_chunk_analyses)
        
        # Create final synthesis message
        user_message_template = prompt_config.get('user_message', '')
        synthesis_message = (
            f"{context_info}"
            f"I've analyzed a large piece of content in chunks. Below are my analyses of each chunk.\n\n"
            f"CHUNK ANALYSES:\n{all_chunk_analyses}\n\n"
            f"Based on these analyses, please provide a final structured analysis "
            f"using the format specified in the system message.\n\n"
            f"{company_context}"
            f"IMPORTANT: Your response MUST include all required structured fields using the specified format."
        )
        
        try:
            # Call OpenAI API for final synthesis
            api_key = os.environ.get("OPENAI_API_KEY", "")
            async_client = AsyncOpenAI(api_key=api_key)
            
            # Begin the API call
            start_time = time.time()
            
            response = await async_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": synthesis_message}
                ],
                temperature=temperature,
                max_tokens=max_tokens
            )
            
            final_analysis = response.choices[0].message.content
            synthesis_tokens = response.usage.total_tokens
            total_tokens += synthesis_tokens
            
            processing_time = time.time() - start_time
            logger.info(f"Final synthesis complete - {synthesis_tokens} tokens used")
            
            return {
                "analysis": final_analysis,
                "model": model,
                "prompt_tokens": 0,  # We don't track separately for chunks
                "completion_tokens": 0,  # We don't track separately for chunks
                "total_tokens": total_tokens,
                "processing_time": processing_time,
                "chunking_used": True,
                "num_chunks": len(chunks)
            }
            
        except Exception as e:
            logger.error(f"Error in final synthesis: {str(e)}")
            return {
                "error": f"Error in final synthesis: {str(e)}",
                "analysis": "Error occurred during content analysis synthesis."
            }

    async def _process_chunks(self, url, content_text, prompt_name, prompt_config, 
                             company_info=None, model=None, retry_attempt=1):
        """
        Process content in chunks when it's too large for a single API call.
        Uses delimited format (field_name ||| value) instead of JSON for more reliable extraction.
        
        Args:
            url: URL being analyzed
            content_text: Large text content to analyze
            prompt_name: Name of the prompt configuration
            prompt_config: Prompt configuration dictionary
            company_info: Optional company context
            model: Model to use (defaults to prompt config)
            retry_attempt: Current retry attempt count
            
        Returns:
            Dictionary containing analysis results
        """
        # Ensure we have a proper model
        if not model:
            model = prompt_config.get('model', 'gpt-4o')
        
        # Calculate appropriate chunk size based on model
        chunk_size = self._get_optimal_chunk_size(model)
        
        logger.info(f"Using content chunking for analysis - content length: {len(content_text)}")
        logger.info(f"Using chunk size of {chunk_size} characters for model: {model}")
        
        # Split content into chunks of appropriate size
        chunks = []
        for i in range(0, len(content_text), chunk_size):
            chunk = content_text[i:i + chunk_size]
            chunks.append(chunk)
            
        # Log chunk statistics
        avg_chunk_size = sum(len(c) for c in chunks) // len(chunks) if chunks else 0
        max_chunk_size = max(len(c) for c in chunks) if chunks else 0
        logger.info(f"Split content into {len(chunks)} chunks: avg={avg_chunk_size} chars, max={max_chunk_size} chars")
        
        # Determine delimiter to use
        delimiter = prompt_config.get('delimiter', '|||')
        
        # Process each chunk
        chunk_results = []
        for i, chunk in enumerate(chunks, 1):
            logger.info(f"Analyzing chunk {i}/{len(chunks)} - size: {len(chunk)} chars")
            
            # Create chunk-specific system message
            chunk_system_message = f"""You are analyzing part {i} of {len(chunks)} of content from {url or 'the document'}.
    Focus on extracting key information from this section only.
    DO NOT make conclusions about the entire content yet.
    Format your observations using field_name {delimiter} value format."""

            if i == 1:
                chunk_system_message += " This is the beginning of the content."
            elif i == len(chunks):
                chunk_system_message += " This is the end of the content."
            
            # Extract prompt configuration
            prompt_system_message = prompt_config.get('system_message', '')
            prompt_user_message = prompt_config.get('user_message', '')
            temperature = prompt_config.get('temperature', 0.3)
            max_tokens = prompt_config.get('max_tokens', 1500)
            
            # Prepare context for the chunk
            page_context = f"URL: {url}\n" if url else ""
            
            # Prepare the company context
            company_context = ""
            if company_info:
                company_context = "Company Information:\n"
                for key, value in company_info.items():
                    company_context += f"- {key.capitalize()}: {value}\n"
            
            # Format user message with the chunk and context
            if "{content}" in prompt_user_message:
                user_message = prompt_user_message.format(
                    content=chunk,
                    company_context=company_context,
                    page_context=page_context
                )
            else:
                user_message = f"{page_context}{prompt_user_message}\n\nCONTENT TO ANALYZE:\n{chunk}\n\n{company_context}"
            
            # Add specific instructions to use delimited format
            user_message += f"\n\nIMPORTANT: Format key insights using field_name {delimiter} value format."
            
            # Call the OpenAI API for this chunk
            try:
                response = await self.async_openai_client.chat.completions.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": chunk_system_message + "\n" + prompt_system_message},
                        {"role": "user", "content": user_message}
                    ],
                    temperature=temperature,
                    max_tokens=max_tokens
                )
                
                # Extract response
                analysis_text = response.choices[0].message.content
                token_usage = response.usage.total_tokens
                
                # Store chunk result
                chunk_results.append({
                    "chunk_index": i,
                    "analysis": analysis_text,
                    "tokens": token_usage
                })
                
                logger.info(f"Chunk {i} analysis complete - {token_usage} tokens used")
                
            except Exception as e:
                logger.error(f"Error processing chunk {i}: {str(e)}")
                chunk_results.append({
                    "chunk_index": i,
                    "analysis": f"Error processing chunk {i}: {str(e)}",
                    "tokens": 0,
                    "error": str(e)
                })
        
        # Now synthesize all chunk results
        logger.info("Synthesizing final analysis from all chunks")
        
        # Create synthesis prompt that asks for delimited format
        synthesis_system_prompt = f"""You are synthesizing a final analysis from multiple content chunks.
    Based on the analyses of {len(chunks)} different sections of the content from {url or 'the document'}, create a comprehensive final analysis.

    MOST IMPORTANTLY: Your response MUST include structured data using the field_name {delimiter} value format for every key insight.

    For example:
    field_name_1 {delimiter} value_1
    field_name_2 {delimiter} value_2

    This format is CRITICAL and must be followed exactly for automated processing.
    """
        
        # Combine all chunk analyses as context for the synthesis
        chunk_contexts = []
        total_tokens = 0
        
        for result in chunk_results:
            chunk_contexts.append(f"--- CHUNK {result['chunk_index']} ANALYSIS ---\n{result['analysis']}")
            total_tokens += result.get('tokens', 0)
        
        combined_analyses = "\n\n".join(chunk_contexts)
        
        # Prepare synthesis user message
        synthesis_user_message = f"""Here is the original prompt:
    {prompt_user_message}

    Below are analyses of {len(chunks)} different chunks of the document from {url or 'the source'}.
    Synthesize these into a comprehensive final analysis.

    {combined_analyses}

    IMPORTANT: Your final analysis MUST include structured data using the field_name {delimiter} value format for every key insight.
    This format MUST be followed exactly as shown in the examples:

    field_name_1 {delimiter} value_1
    field_name_2 {delimiter} value_2

    Include all the fields requested in the original prompt."""
        
        # Call OpenAI API for synthesis
        try:
            response = await self.async_openai_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": synthesis_system_prompt},
                    {"role": "user", "content": synthesis_user_message}
                ],
                temperature=0.3,  # Lower temperature for more accurate structured data generation
                max_tokens=2000   # Allow more tokens for synthesis
            )
            
            # Extract synthesis result
            synthesis_text = response.choices[0].message.content
            synthesis_tokens = response.usage.total_tokens
            
            logger.info(f"Final synthesis complete - {synthesis_tokens} tokens used")
            
            # Extract structured data using delimited format extraction
            structured_data = self._extract_delimited_fields(synthesis_text, delimiter, prompt_name)
            
            # If no structured data found, try once more with explicit instructions
            if not structured_data and retry_attempt < 3:
                logger.warning(f"No structured data found for '{prompt_name}' (attempt {retry_attempt}/3)")
                
                # Sleep briefly before retrying
                await asyncio.sleep(1)
                
                # Retry with more explicit instructions
                return await self._process_chunks(
                    url=url,
                    content_text=content_text,
                    prompt_name=prompt_name,
                    prompt_config=prompt_config,
                    company_info=company_info,
                    model=model,
                    retry_attempt=retry_attempt + 1
                )
            
            # Prepare final result
            result = {
                "analysis": synthesis_text,
                "structured_data": {prompt_name: structured_data} if structured_data else {},
                "model": model,
                "total_tokens": total_tokens + synthesis_tokens,
                "api_tokens": total_tokens + synthesis_tokens,  # For consistency with non-chunked processing
                "prompt_tokens": 0,  # Not available in chunked processing
                "completion_tokens": 0,  # Not available in chunked processing
                "processing_time": 0,  # Not tracked precisely in chunked processing
                "chunks": len(chunks)
            }
            
            return result
            
        except Exception as e:
            logger.error(f"Error in synthesis: {str(e)}")
            return {
                "error": f"Synthesis error: {str(e)}",
                "analysis": "Error synthesizing results from chunks.",
                "model": model,
                "total_tokens": total_tokens,
                "chunks": len(chunks)
            }

    def _extract_delimited_fields(self, text, delimiter="|||", prompt_name="unknown"):
        """
        Extract fields from text using a delimiter format (field_name delimiter value).
        
        Args:
            text: Text containing delimited fields
            delimiter: The delimiter string (default: "|||")
            prompt_name: Name of the prompt for logging
            
        Returns:
            Dictionary of extracted fields
        """
        if not text:
            return {}
        
        import re
        
        # Prepare the delimiter for regex (escape special characters)
        esc_delimiter = re.escape(delimiter)
        
        # Multiple patterns to try, from most to least specific
        patterns = [
            # Standard pattern with word boundary: field_name ||| value
            rf"([a-zA-Z0-9_]+)\s*{esc_delimiter}\s*(.+?)(?:\n|$)",
            
            # Alternative with possible prefix
            rf"[^a-zA-Z0-9_]*([a-zA-Z0-9_]+)\s*{esc_delimiter}\s*(.+?)(?:\n|$)",
            
            # Allow for possible formatting/bullets
            rf"[-*•]?\s*([a-zA-Z0-9_]+)\s*{esc_delimiter}\s*(.+?)(?:\n|$)"
        ]
        
        fields = {}
        matches_found = False
        
        # Try each pattern until we find something
        for pattern in patterns:
            matches = re.findall(pattern, text)
            if matches:
                matches_found = True
                for field_name, field_value in matches:
                    if field_name.strip() and field_value.strip():
                        # Clean up the field value
                        clean_value = field_value.strip()
                        fields[field_name.strip()] = self._convert_value_type(clean_value)
        
        # Special case: if no matches but there's "|||" in the text, try raw line extraction
        if not matches_found and delimiter in text:
            for line in text.split('\n'):
                if delimiter in line:
                    parts = line.split(delimiter, 1)
                    if len(parts) == 2:
                        field_name = parts[0].strip()
                        field_value = parts[1].strip()
                        
                        # Clean up field name - remove any non-alphanumeric prefix
                        field_name = re.sub(r'^[^a-zA-Z0-9_]*', '', field_name)
                        
                        if field_name and field_value:
                            fields[field_name] = self._convert_value_type(field_value)
        
        if fields:
            logger.info(f"Successfully extracted {len(fields)} delimited fields for '{prompt_name}'")
        else:
            logger.warning(f"No delimited fields found for '{prompt_name}'")
        
        return fields

    def _convert_value_type(self, value):
        """Convert string values to appropriate types."""
        # Remove quotes if present
        if (value.startswith('"') and value.endswith('"')) or \
           (value.startswith("'") and value.endswith("'")):
            value = value[1:-1]
        
        value = value.strip()
        
        try:
            # Check if it's a numerical value
            if value.isdigit():
                return int(value)
            # Check if it's a float
            elif value.replace('.', '', 1).isdigit() and value.count('.') == 1:
                return float(value)
            # Check if it's a boolean
            elif value.lower() in ["true", "yes"]:
                return True
            elif value.lower() in ["false", "no"]:
                return False
            # Check if it's a list (comma separated values)
            elif "," in value and not "[" in value:
                return [item.strip() for item in value.split(",")]
            # Check if it's "None"
            elif value.lower() == "none":
                return None
            # Otherwise keep as string
            else:
                return value
        except Exception:
            # If conversion fails, keep the original string
            return value



    def _smart_chunk_content(self, content, max_size=None, model=None):
        """
        Intelligently chunk content into logical sections based on structure and model capabilities.
        
        Args:
            content: Content text to chunk
            max_size: Maximum size per chunk (characters). If None, determined from model.
            model: Model name to determine optimal chunk size
            
        Returns:
            List of content chunks
        """
        # Determine optimal chunk size based on model
        if max_size is None:
            if model:
                max_size = self._get_optimal_chunk_size(model)
            else:
                # Default to a reasonable chunk size for newer models
                max_size = 40000
        
        logger.info(f"Using chunk size of {max_size} characters for model: {model if model else 'default'}")
        
        # If content fits in one chunk, return it directly
        if len(content) <= max_size:
            return [content]
        
        # Find logical break points using headers, section breaks, etc.
        import re
        
        # Split at major section breaks first (paragraphs, headers)
        # Look for double newlines, HTML headers, or other section markers
        section_pattern = r'((?:\r?\n){2,}|(?:<h[1-6]>)|(?:<\/h[1-6]>)|(?:^#{1,6}\s+.*$))'
        sections = re.split(section_pattern, content, flags=re.MULTILINE)
        
        # Process sections into chunks
        chunks = []
        current_chunk = ""
        
        for section in sections:
            # If adding this section keeps us under the max size, add it
            if len(current_chunk) + len(section) <= max_size:
                current_chunk += section
            else:
                # If current section is very large, we need to split it further
                if len(section) > max_size:
                    # If we have accumulated content, add it as a chunk first
                    if current_chunk:
                        chunks.append(current_chunk)
                        current_chunk = ""
                    
                    # For long sections, split at paragraph level
                    paragraphs = re.split(r'(\r?\n)', section)
                    for paragraph in paragraphs:
                        if len(current_chunk) + len(paragraph) <= max_size:
                            current_chunk += paragraph
                        else:
                            # If paragraph is still too large, split by sentences
                            if len(paragraph) > max_size:
                                if current_chunk:
                                    chunks.append(current_chunk)
                                    current_chunk = ""
                                    
                                # Split by sentences (try to keep sentences together)
                                sentence_pattern = r'([.!?]+\s+)'
                                sentences = re.split(sentence_pattern, paragraph)
                                for sentence in sentences:
                                    if len(current_chunk) + len(sentence) <= max_size:
                                        current_chunk += sentence
                                    else:
                                        if current_chunk:
                                            chunks.append(current_chunk)
                                        
                                        # If a single sentence is too large, split arbitrarily
                                        if len(sentence) > max_size:
                                            # Use a slightly smaller size for arbitrary splits
                                            # to ensure we don't exceed max_size
                                            safe_size = max(100, max_size - 100)
                                            for i in range(0, len(sentence), safe_size):
                                                chunks.append(sentence[i:i+safe_size])
                                            current_chunk = ""
                                        else:
                                            current_chunk = sentence
                            else:
                                # Add accumulated content as a chunk
                                if current_chunk:
                                    chunks.append(current_chunk)
                                current_chunk = paragraph
                else:
                    # Add accumulated content as a chunk
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = section
        
        # Add the last chunk if there's anything left
        if current_chunk:
            chunks.append(current_chunk)
        
        # Log chunk statistics
        chunk_sizes = [len(chunk) for chunk in chunks]
        avg_size = sum(chunk_sizes) / len(chunks) if chunks else 0
        max_chunk = max(chunk_sizes) if chunks else 0
        
        logger.info(f"Split content into {len(chunks)} chunks: avg={int(avg_size)} chars, max={max_chunk} chars")
        
        return chunks

    def _get_optimal_chunk_size(self, model_name):
        """
        Get optimal chunk size for a given model based on its context window.
        
        Args:
            model_name: Name of the model (e.g., 'o3-mini', 'gpt-4o', 'o1')
        
        Returns:
            Optimal chunk size in characters
        """
        model_name = model_name.lower()
        
        # Set reasonable limits based on model capabilities (in characters, not tokens)
        # Using conservative estimates (roughly 4 chars per token)
        
        # Newer models (as of 2025)
        if 'o3-mini' in model_name:
            # o3-mini: 200K context window, 100K max output
            # Use ~60K chars per chunk (conservative to leave room for prompt and output)
            return 60000
        elif 'o1' in model_name:
            # o1: 200K context window, 100K max output
            return 60000
        
        # Previous models
        elif 'gpt-4o' in model_name:
            # GPT-4o: 128K context window, 16K max output
            return 40000
        elif 'gpt-4-turbo' in model_name:
            # GPT-4 Turbo: 128K context window
            return 40000
        elif 'gpt-4-32k' in model_name:
            # GPT-4-32k: 32K context window
            return 20000
        elif 'gpt-4' in model_name:
            # Standard GPT-4: 8K context window
            return 12000
        elif 'gpt-3.5-turbo-16k' in model_name:
            # GPT-3.5-16k: 16K context window
            return 10000
        elif 'gpt-3.5' in model_name:
            # Standard GPT-3.5: 4K context window
            return 6000
        else:
            # Default for unknown models - use a moderate size
            # This will work for most use cases but may not be optimal
            return 15000


def preprocess_content(content_text):
    """
    Preprocess content text to improve analysis quality.
    
    Args:
        content_text: Raw content text from scraping
        
    Returns:
        Preprocessed content
    """
    if not content_text:
        return ""
        
    # Remove excessive whitespace
    content_text = re.sub(r'\s+', ' ', content_text)
    
    # Remove repetitive elements (common in poorly scraped pages)
    lines = content_text.split('\n')
    unique_lines = []
    seen = set()
    
    for line in lines:
        line_stripped = line.strip()
        if line_stripped and line_stripped not in seen:
            seen.add(line_stripped)
            unique_lines.append(line)
            
    # Join lines back together
    content_text = '\n'.join(unique_lines)
    
    # Add markers to clearly delineate content
    content_text = f"--- WEB CONTENT START ---\n{content_text}\n--- WEB CONTENT END ---"
    
    return content_text



class URLProcessor:
    """
    Handles URL batch processing with URL-specific company context.
    """
    
    def __init__(self, analyzer: ContentAnalyzer, db_manager: DatabaseManager):
        """Initialize the URL processor."""
        self.analyzer = analyzer
        self.db = db_manager
    
    async def _process_urls_async(self, job_id: str, url_data_list: List[Dict[str, Any]], prompt_names: List[str]):
        """
        Process URLs asynchronously in batches with controlled concurrency.
        """
        total_urls = len(url_data_list)
        processed = 0
        errors = 0
        
        # Determine optimal batch size and concurrency
        batch_size = int(os.environ.get("BATCH_SIZE", "10"))
        max_concurrency = int(os.environ.get("MAX_CONCURRENCY", "5"))
        
        logger.info(f"Processing job {job_id} with {total_urls} URLs in batches of {batch_size} with max concurrency {max_concurrency}")
        
        # Split into batches
        batches = [url_data_list[i:i + batch_size] for i in range(0, len(url_data_list), batch_size)]
        
        # Semaphore to control concurrency
        semaphore = asyncio.Semaphore(max_concurrency)
        
        async def process_url_with_semaphore(url_data):
            """Process a single URL with semaphore control"""
            async with semaphore:
                url = url_data.get('url')
                company_info = url_data.get('company_info')
                content_type = url_data.get('content_type', 'html')
                force_browser = url_data.get('force_browser', False)
                
                try:
                    result = await self.analyzer.process_url_async(
                        url=url,
                        prompt_names=prompt_names,
                        company_info=company_info,
                        content_type=content_type, 
                        force_browser=force_browser,
                        job_id=job_id 
                    )
                    result['job_id'] = job_id
                    return result
                except Exception as e:
                    logger.error(f"Error processing URL {url}: {str(e)}")
                    return {
                        "url": url,
                        "job_id": job_id,
                        "status": "error",
                        "error": str(e),
                        "created_at": datetime.now().isoformat()
                    }
        
        # Process each batch
        for i, batch in enumerate(batches):
            logger.info(f"Processing batch {i+1}/{len(batches)} for job {job_id}")
            
            # Create tasks for all URLs in the batch
            tasks = [process_url_with_semaphore(url_data) for url_data in batch if url_data.get('url')]
            
            # Process the batch
            batch_results = await asyncio.gather(*tasks, return_exceptions=True)
            
            # Save results and update counters
            for result in batch_results:
                if isinstance(result, Exception):
                    logger.error(f"Exception in URL processing: {str(result)}")
                    errors += 1
                    continue
                
                # Save result to database
                self.db.save_result(result)
                
                # Update counters
                processed += 1
                if result.get('status') != 'success':
                    errors += 1
            
            # Update job status after each batch
            self.db.update_job_status(
                job_id=job_id,
                processed_urls=processed,
                error_count=errors
            )
            
            logger.info(f"Completed batch {i+1}/{len(batches)} for job {job_id} - {processed}/{total_urls} processed, {errors} errors")
        
        # Mark job as completed
        final_status = "completed" if errors == 0 else "completed_with_errors"
        self.db.update_job_status(
            job_id=job_id,
            status=final_status,
            processed_urls=processed,
            error_count=errors,
            completed_at=datetime.now().isoformat()
        )
        
        logger.info(f"Finished job {job_id}: {processed}/{total_urls} URLs processed, {errors} errors")
    
    def process_url_batch(self, job_id: str, url_data_list: List[Dict[str, Any]], prompt_names: List[str]):
        """
        Process a batch of URLs with their specific company contexts.
        Wrapper for the async version.
        
        Args:
            job_id: The job identifier
            url_data_list: List of dictionaries containing URL and its associated company context
                Each dictionary should have at least a 'url' key and optionally a 'company_info' key
            prompt_names: List of prompt configuration names to use
        """
        # Update job status to running
        self.db.update_job_status(
            job_id=job_id,
            status="running",
            total_urls=len(url_data_list),
            processed_urls=0,
            error_count=0
        )
        
        # Run the async version in a new event loop
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            loop.run_until_complete(
                self._process_urls_async(job_id, url_data_list, prompt_names)
            )
        finally:
            loop.close()
